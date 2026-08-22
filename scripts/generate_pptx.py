import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    COLOR_BG = RGBColor(7, 9, 14)          # Dark #07090E
    COLOR_CARD = RGBColor(14, 18, 27)      # Card #0E121B
    COLOR_CYAN = RGBColor(0, 229, 255)     # Cyan #00E5FF
    COLOR_GREEN = RGBColor(0, 230, 118)    # Green #00E676
    COLOR_ORANGE = RGBColor(255, 145, 0)   # Orange #FF9100
    COLOR_RED = RGBColor(255, 61, 0)       # Red #FF3D00
    COLOR_TEXT_MAIN = RGBColor(240, 244, 248) # Crisp #F0F4F8
    COLOR_TEXT_MUTED = RGBColor(138, 153, 173) # Muted #8A99AD
    COLOR_BORDER = RGBColor(30, 38, 56)    # Border #1E2638

    def set_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, tag_text, title_text, subtitle_text):
        # Tag
        tb_tag = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_tag = tb_tag.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = tag_text.upper()
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.name = "Arial"
        p_tag.font.color.rgb = COLOR_CYAN

        # Title
        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.8))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.name = "Arial"
        p_title.font.color.rgb = COLOR_TEXT_MAIN

        # Subtitle
        if subtitle_text:
            tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(11.7), Inches(0.5))
            tf_sub = tb_sub.text_frame
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle_text
            p_sub.font.size = Pt(14)
            p_sub.font.name = "Arial"
            p_sub.font.color.rgb = COLOR_TEXT_MUTED

    def add_card(slide, left, top, width, height, border_color=None):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
        else:
            card.line.color.rgb = COLOR_BORDER
            card.line.width = Pt(1)
        return card

    def add_footer(slide, current_idx):
        tb_ft = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.733), Inches(0.4))
        tf_ft = tb_ft.text_frame
        p_ft = tf_ft.paragraphs[0]
        p_ft.text = f"Thervo — Yuva Yodha Tech Challenge (Smart Buildings)                          Slide {current_idx} / 10"
        p_ft.font.size = Pt(10)
        p_ft.font.name = "Arial"
        p_ft.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 1: TITLE ====================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1)
    
    tb1_tag = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(0.4))
    p = tb1_tag.text_frame.paragraphs[0]
    p.text = "YUVA YODHA TECH CHALLENGE  •  SMART BUILDINGS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    tb1_title = s1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.4))
    p = tb1_title.text_frame.paragraphs[0]
    p.text = "Thervo"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN

    tb1_sub = s1.shapes.add_textbox(Inches(1.0), Inches(3.6), Inches(11.3), Inches(0.8))
    p = tb1_sub.text_frame.paragraphs[0]
    p.text = "Predictive Cooling for a More Efficient Digital World"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    tb1_line = s1.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.3), Inches(0.6))
    p = tb1_line.text_frame.paragraphs[0]
    p.text = "AI-driven, sensor-free thermal intelligence for data centers"
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_TEXT_MUTED

    add_card(s1, Inches(1.0), Inches(5.4), Inches(4.5), Inches(0.7), COLOR_CYAN)
    tb1_badge1 = s1.shapes.add_textbox(Inches(1.1), Inches(5.45), Inches(4.3), Inches(0.6))
    p = tb1_badge1.text_frame.paragraphs[0]
    p.text = "[ Predictive Thermal Risk Engine ]"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.alignment = PP_ALIGN.CENTER

    add_card(s1, Inches(5.8), Inches(5.4), Inches(4.5), Inches(0.7), COLOR_GREEN)
    tb1_badge2 = s1.shapes.add_textbox(Inches(5.9), Inches(5.45), Inches(4.3), Inches(0.6))
    p = tb1_badge2.text_frame.paragraphs[0]
    p.text = "[ Zero Hardware Sensors Required ]"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    p.alignment = PP_ALIGN.CENTER

    add_footer(s1, 1)

    # ==================== SLIDE 2: THE PROBLEM ====================
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2)
    add_header(s2, "Problem Analysis", "Cooling is becoming a hidden cost of our digital future", "Traditional data center cooling cannot keep up with high-density AI workloads.")

    add_card(s2, Inches(0.8), Inches(2.1), Inches(6.0), Inches(4.5))
    tb2_left = s2.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.6), Inches(4.1))
    tf2_left = tb2_left.text_frame
    tf2_left.word_wrap = True

    bullets2 = [
        "Expanding Energy Demand: AI and HPC create unprecedented power density per server rack.",
        "Major Energy Share: Cooling represents up to 40% of overall data center electricity usage.",
        "Reactive Model: Conventional cooling triggers only AFTER heat has already accumulated.",
        "Blanket Inefficiency: Entire aisles flooded with chilled air even when single racks run hot.",
        "Delayed Sensor Feedback: Physical probes register heat after thermal stress has developed."
    ]
    for idx, b in enumerate(bullets2):
        p = tf2_left.paragraphs[0] if idx == 0 else tf2_left.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(10)

    add_card(s2, Inches(7.1), Inches(2.1), Inches(5.4), Inches(4.5), COLOR_RED)
    tb2_rtitle = s2.shapes.add_textbox(Inches(7.3), Inches(2.3), Inches(5.0), Inches(0.5))
    p = tb2_rtitle.text_frame.paragraphs[0]
    p.text = "TODAY'S REACTIVE COOLING LOOP"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED

    steps2 = [
        ("1. Workload Spikes", "CPU/GPU load rises"),
        ("2. Heat Develops", "Physical heat accumulates"),
        ("3. Sensor Detects", "Ambient probe triggers"),
        ("4. Cooling Reacts", "Floods entire zone")
    ]
    for idx, (title, desc) in enumerate(steps2):
        top_y = 2.9 + (idx * 0.85)
        add_card(s2, Inches(7.4), Inches(top_y), Inches(4.8), Inches(0.65))
        tb_st = s2.shapes.add_textbox(Inches(7.5), Inches(top_y + 0.05), Inches(4.6), Inches(0.55))
        p = tb_st.text_frame.paragraphs[0]
        p.text = f"{title} — {desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_CYAN if idx == 3 else COLOR_TEXT_MAIN

    add_footer(s2, 2)

    # ==================== SLIDE 3: WHY IT MATTERS ====================
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3)
    add_header(s3, "Strategic Relevance", "We shouldn't have to cool everything to protect anything", "Precision and foresight are critical for modern data infrastructure sustainability.")

    add_card(s3, Inches(0.8), Inches(2.1), Inches(6.0), Inches(4.5))
    tb3_left = s3.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.6), Inches(2.5))
    tf3_left = tb3_left.text_frame
    tf3_left.word_wrap = True

    bullets3 = [
        "Relentless Compute Growth: Higher density makes blanket cooling cost-prohibitive.",
        "Operational Inflation: Blanket cooling continuously inflates facility OpEx.",
        "Thermal Stress Risks: Delayed cooling exposes hardware to throttling & degradation.",
        "Predictive Shift Needed: Operators require rack-level foresight before heat spikes."
    ]
    for idx, b in enumerate(bullets3):
        p = tf3_left.paragraphs[0] if idx == 0 else tf3_left.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(10)

    # Callout
    add_card(s3, Inches(1.0), Inches(4.9), Inches(5.6), Inches(1.4), COLOR_CYAN)
    tb3_call = s3.shapes.add_textbox(Inches(1.1), Inches(5.0), Inches(5.4), Inches(1.2))
    tf3_call = tb3_call.text_frame
    tf3_call.word_wrap = True
    p = tf3_call.paragraphs[0]
    p.text = "\"Thervo asks a simple question: What if we could predict where cooling will be needed before the heat arrives?\""
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    # Right Card
    add_card(s3, Inches(7.1), Inches(2.1), Inches(5.4), Inches(4.5), COLOR_GREEN)
    tb3_rtitle = s3.shapes.add_textbox(Inches(7.3), Inches(2.3), Inches(5.0), Inches(0.5))
    p = tb3_rtitle.text_frame.paragraphs[0]
    p.text = "TARGETED vs BLANKET COOLING"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    add_card(s3, Inches(7.4), Inches(3.0), Inches(4.8), Inches(1.5), COLOR_RED)
    tb3_b1 = s3.shapes.add_textbox(Inches(7.5), Inches(3.1), Inches(4.6), Inches(1.3))
    tf3_b1 = tb3_b1.text_frame
    tf3_b1.word_wrap = True
    p = tf3_b1.paragraphs[0]
    p.text = "TRADITIONAL BLANKET COOLING"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED
    p2 = tf3_b1.add_paragraph()
    p2.text = "Entire server hall chilled uniformly regardless of individual rack utilization."
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_TEXT_MUTED

    add_card(s3, Inches(7.4), Inches(4.7), Inches(4.8), Inches(1.6), COLOR_GREEN)
    tb3_b2 = s3.shapes.add_textbox(Inches(7.5), Inches(4.8), Inches(4.6), Inches(1.4))
    tf3_b2 = tb3_b2.text_frame
    tf3_b2.word_wrap = True
    p = tf3_b2.paragraphs[0]
    p.text = "THERVO TARGETED MICRO-ZONE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    p2 = tf3_b2.add_paragraph()
    p2.text = "Cooling airflow directed specifically to identified high-risk racks pre-emptively."
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_TEXT_MUTED

    add_footer(s3, 3)

    # ==================== SLIDE 4: THE SOLUTION ====================
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4)
    add_header(s4, "Core Solution", "Meet Thervo", "AI-driven predictive cooling operating in 3 simple steps.")

    col_w = Inches(3.7)
    steps4 = [
        ("STEP 01", "Understand Workload", "Ingests server performance telemetry directly from host systems: CPU, GPU, Memory, Disk I/O, Network I/O.", COLOR_CYAN),
        ("STEP 02", "Predict Thermal Risk", "Machine learning models convert real-time workload behavior into pre-emptive rack-level thermal risk scores.", COLOR_GREEN),
        ("STEP 03", "Act Before Problem", "Cooling airflow is targeted directly toward specific racks where thermal risk is actively developing.", COLOR_ORANGE)
    ]
    for idx, (st_tag, st_title, st_desc, st_color) in enumerate(steps4):
        left_x = Inches(0.8) + idx * Inches(4.0)
        add_card(s4, left_x, Inches(2.1), col_w, Inches(3.2), st_color)
        tb_step = s4.shapes.add_textbox(left_x + Inches(0.2), Inches(2.3), col_w - Inches(0.4), Inches(2.8))
        tf_step = tb_step.text_frame
        tf_step.word_wrap = True
        p = tf_step.paragraphs[0]
        p.text = st_tag
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = st_color
        p2 = tf_step.add_paragraph()
        p2.text = st_title
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT_MAIN
        p2.space_before = Pt(6)
        p3 = tf_step.add_paragraph()
        p3.text = st_desc
        p3.font.size = Pt(12)
        p3.font.color.rgb = COLOR_TEXT_MUTED
        p3.space_before = Pt(10)

    # Key Differentiator
    add_card(s4, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.9), COLOR_GREEN)
    tb4_diff = s4.shapes.add_textbox(Inches(1.0), Inches(5.75), Inches(11.3), Inches(0.6))
    p = tb4_diff.text_frame.paragraphs[0]
    p.text = "KEY DIFFERENTIATOR: NO DEDICATED TEMPERATURE SENSORS REQUIRED"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    p.alignment = PP_ALIGN.CENTER

    add_footer(s4, 4)

    # ==================== SLIDE 5: HOW THERVO WORKS ====================
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5)
    add_header(s5, "Technical Architecture", "From workload signals to thermal intelligence", "Data pipeline and algorithmic foundation driving pre-emptive cooling decisions.")

    # Flow Box
    add_card(s5, Inches(0.8), Inches(2.1), Inches(11.7), Inches(1.3), COLOR_CYAN)
    tb5_flow = s5.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.1))
    tf5_flow = tb5_flow.text_frame
    tf5_flow.word_wrap = True
    p = tf5_flow.paragraphs[0]
    p.text = "SYSTEM PIPELINE FLOW:"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p2 = tf5_flow.add_paragraph()
    p2.text = "Server Workload Telemetry ➔ Feature Processing ➔ 15-Feature Vector ➔ XGBoost Model + Graph Heat Propagation ➔ Composite Rack Risk ➔ Cooling Orchestration"
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_MAIN
    p2.space_before = Pt(6)

    # Two columns
    add_card(s5, Inches(0.8), Inches(3.6), Inches(5.7), Inches(3.0))
    tb5_l = s5.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(5.3), Inches(2.6))
    tf5_l = tb5_l.text_frame
    tf5_l.word_wrap = True
    p = tf5_l.paragraphs[0]
    p.text = "PRIMARY ENGINE COMPONENTS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    b_l = [
        "Workload Heat Proxy: CPU and GPU utilization metrics act as heat proxies.",
        "Graph Thermal Context: Neighboring rack influence modeled via graph propagation.",
        "XGBoost Risk Engine: Fast, deterministic model yields primary risk scores."
    ]
    for b in b_l:
        p = tf5_l.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT_MAIN
        p.space_before = Pt(8)

    add_card(s5, Inches(6.8), Inches(3.6), Inches(5.7), Inches(3.0))
    tb5_r = s5.shapes.add_textbox(Inches(7.0), Inches(3.8), Inches(5.3), Inches(2.6))
    tf5_r = tb5_r.text_frame
    tf5_r.word_wrap = True
    p = tf5_r.paragraphs[0]
    p.text = "MODEL SPECIFICATIONS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    b_r = [
        "Composite Fusion: Fuses single-rack XGBoost signal with spatial graph context.",
        "Normalized Bounding: Outputs strictly bounded within [0.0, 1.0] domain.",
        "Zero Hardware Overhead: Lightweight inference runs on existing compute nodes."
    ]
    for b in b_r:
        p = tf5_r.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT_MAIN
        p.space_before = Pt(8)

    add_footer(s5, 5)

    # ==================== SLIDE 6: USER JOURNEY ====================
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6)
    add_header(s6, "Product Experience", "What an operator sees", "Mission control workflow for real-time facility visibility and pre-emptive control.")

    add_card(s6, Inches(0.8), Inches(2.1), Inches(6.0), Inches(4.5))
    tb6_left = s6.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.6), Inches(4.1))
    tf6_left = tb6_left.text_frame
    tf6_left.word_wrap = True

    journey_steps = [
        "1. Operator opens Thervo Mission Control dashboard.",
        "2. Views real-time rack-level thermal risk heat grid.",
        "3. Identifies emerging high-risk racks prior to heat buildup.",
        "4. Receives predictive alerts before critical thresholds.",
        "5. Automated cooling triggers engage when risk thresholds breach.",
        "6. Operator retains real-time manual override control.",
        "7. Risk history and energy telemetry log long-term performance."
    ]
    for idx, j in enumerate(journey_steps):
        p = tf6_left.paragraphs[0] if idx == 0 else tf6_left.add_paragraph()
        p.text = j
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_CYAN if idx in [3, 4] else COLOR_TEXT_MAIN
        p.space_after = Pt(6)

    # Dashboard Mockup Card
    add_card(s6, Inches(7.1), Inches(2.1), Inches(5.4), Inches(4.5), COLOR_CYAN)
    tb6_dash = s6.shapes.add_textbox(Inches(7.3), Inches(2.3), Inches(5.0), Inches(4.1))
    tf6_dash = tb6_dash.text_frame
    tf6_dash.word_wrap = True
    p = tf6_dash.paragraphs[0]
    p.text = "MISSION CONTROL DASHBOARD MOCKUP"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p2 = tf6_dash.add_paragraph()
    p2.text = "Live Grid Status:\n• RACK-01: Risk 0.14 [NOMINAL]\n• RACK-02: Risk 0.42 [COOLING ACTIVE]\n• RACK-03: Risk 0.78 [WARNING]\n• RACK-04: Risk 0.91 [CRITICAL ALERT]\n• RACK-05: Risk 0.22 [NOMINAL]\n• RACK-06: Risk 0.18 [NOMINAL]"
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_TEXT_MAIN
    p2.space_before = Pt(8)

    p3 = tf6_dash.add_paragraph()
    p3.text = "PREDICTIVE ALERT: Rack-04 thermal risk spiking -> TRIGGERING COOLING"
    p3.font.size = Pt(11)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_RED
    p3.space_before = Pt(12)

    add_footer(s6, 6)

    # ==================== SLIDE 7: DIFFERENTIATION ====================
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7)
    add_header(s7, "Competitive Differentiation", "Predictive. Sensor-free. Rack-level.", "Reinventing cooling strategy from reactive zones to workload-aware micro-zones.")

    # Table
    table_shape = s7.shapes.add_table(5, 3, Inches(0.8), Inches(2.1), Inches(11.7), Inches(2.4))
    table = table_shape.table
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(4.35)
    table.columns[2].width = Inches(4.35)

    headers = ["Dimension", "Traditional Cooling Systems", "Thervo Thermal Intelligence"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_CARD
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLOR_CYAN if i == 2 else COLOR_TEXT_MUTED

    rows_data = [
        ("Operational Mode", "Reactive (Responds post-heat)", "Predictive (Acts pre-heat)"),
        ("Cooling Granularity", "Broad / Zone-Level Flooding", "Micro-Zone / Rack-Level"),
        ("Hardware Dependency", "Requires Physical Sensors", "Uses Existing Telemetry (Sensor-Free)"),
        ("Decision Logic", "Static Temperature Alarms", "AI-Driven Risk Scoring")
    ]
    for row_idx, r in enumerate(rows_data):
        for col_idx, val in enumerate(r):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_BG
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = COLOR_CYAN if col_idx == 2 else COLOR_TEXT_MAIN

    # 3 Cards below
    c_w = Inches(3.7)
    innovations = [
        ("1. Sensor-Free Inference", "Utilizes native OS server performance signals, eliminating hardware deployment & maintenance costs.", COLOR_CYAN),
        ("2. Workload-Aware Foresight", "Anticipates thermal generation directly from CPU/GPU compute spikes before heat builds up.", COLOR_GREEN),
        ("3. Graph Neighbor Model", "Models thermal bleed and heat dissipation between adjacent server racks in the aisle.", COLOR_ORANGE)
    ]
    for idx, (title, desc, color) in enumerate(innovations):
        left_x = Inches(0.8) + idx * Inches(4.0)
        add_card(s7, left_x, Inches(4.8), c_w, Inches(1.8), color)
        tb_in = s7.shapes.add_textbox(left_x + Inches(0.2), Inches(4.9), c_w - Inches(0.4), Inches(1.6))
        tf_in = tb_in.text_frame
        tf_in.word_wrap = True
        p = tf_in.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        p2 = tf_in.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED
        p2.space_before = Pt(6)

    add_footer(s7, 7)

    # ==================== SLIDE 8: EXPECTED IMPACT ====================
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8)
    add_header(s8, "Expected Impact", "Less wasted cooling. Earlier intervention. Smarter infrastructure.", "Value creation across energy, operational cost, reliability, and sustainability.")

    impacts = [
        ("ENERGY", "Targeted Cooling", "Cuts unnecessary cooling by delivering airflow strictly where heat is predicted to accumulate.", COLOR_CYAN),
        ("COST", "OpEx & CapEx Efficiency", "Lowers cooling electricity OpEx and avoids capital expense for physical sensor hardware.", COLOR_GREEN),
        ("RELIABILITY", "Early Hazard Mitigation", "Identifies emerging thermal risks earlier, reducing hardware stress and thermal throttling.", COLOR_ORANGE),
        ("SUSTAINABILITY", "Lower PUE Profile", "Supports greener, lower PUE operations for energy-intensive AI digital infrastructure.", COLOR_CYAN)
    ]
    i_w = Inches(2.7)
    for idx, (tag, title, desc, color) in enumerate(impacts):
        left_x = Inches(0.8) + idx * Inches(2.95)
        add_card(s8, left_x, Inches(2.1), i_w, Inches(3.2), color)
        tb_i = s8.shapes.add_textbox(left_x + Inches(0.15), Inches(2.25), i_w - Inches(0.3), Inches(2.9))
        tf_i = tb_i.text_frame
        tf_i.word_wrap = True
        p = tf_i.paragraphs[0]
        p.text = tag
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        p2 = tf_i.add_paragraph()
        p2.text = title
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT_MAIN
        p2.space_before = Pt(4)
        p3 = tf_i.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = COLOR_TEXT_MUTED
        p3.space_before = Pt(8)

    add_card(s8, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.1))
    tb8_bot = s8.shapes.add_textbox(Inches(1.0), Inches(5.6), Inches(11.3), Inches(0.9))
    tf8_bot = tb8_bot.text_frame
    tf8_bot.word_wrap = True
    p = tf8_bot.paragraphs[0]
    p.text = "KEY BENEFICIARIES: Data-center operators | Infrastructure/SRE teams | Facilities managers | AI cloud providers"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p2 = tf8_bot.add_paragraph()
    p2.text = "Note: Impact metrics represent targeted model goals and simulation benchmarks. Real-world savings vary by facility architecture."
    p2.font.size = Pt(10)
    p2.font.color.rgb = COLOR_TEXT_MUTED
    p2.space_before = Pt(4)

    add_footer(s8, 8)

    # ==================== SLIDE 9: ROADMAP ====================
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9)
    add_header(s9, "Execution Plan", "From validated system to production-scale cooling intelligence", "Clear distinction between completed validation and future deployment phases.")

    phases = [
        ("PHASE 0 • COMPLETED", "Simulation & Pipeline Validation", "Interactive simulator, ML pipeline, synthetic & hardware telemetry validation", COLOR_GREEN),
        ("PHASE 1", "Pilot Deployment", "Live telemetry integration, single-floor deployment, baseline energy comparison", COLOR_CYAN),
        ("PHASE 2", "Production Scale Expansion", "Multi-floor deployment, DCIM integration, production API & role-based access control", COLOR_TEXT_MUTED),
        ("PHASE 3", "Advanced Intelligence Layer", "Temporal forecasting, anomaly detection, digital twin simulation capabilities", COLOR_TEXT_MUTED),
        ("PHASE 4", "Global Enterprise Platform", "Multi-facility orchestration, enterprise SaaS offering, global partner ecosystem", COLOR_TEXT_MUTED)
    ]
    for idx, (p_tag, p_title, p_desc, p_color) in enumerate(phases):
        top_y = 2.1 + (idx * 0.9)
        border_c = COLOR_GREEN if idx == 0 else (COLOR_CYAN if idx == 1 else None)
        add_card(s9, Inches(0.8), Inches(top_y), Inches(11.7), Inches(0.75), border_c)
        tb_p = s9.shapes.add_textbox(Inches(1.0), Inches(top_y + 0.1), Inches(11.3), Inches(0.55))
        tf_p = tb_p.text_frame
        tf_p.word_wrap = True
        p = tf_p.paragraphs[0]
        p.text = f"{p_tag}  |  {p_title}  —  {p_desc}"
        p.font.size = Pt(12)
        p.font.bold = (idx <= 1)
        p.font.color.rgb = p_color

    add_footer(s9, 9)

    # ==================== SLIDE 10: TEAM + CLOSING ====================
    s10 = prs.slides.add_slide(blank_layout)
    set_bg(s10)
    add_header(s10, "Team & Vision", "Building a smarter way to cool the digital world", "Engineered for real-world data center sustainability.")

    team_members = [
        ("AI & MACHINE LEARNING", "[ Team Member 1 ]", "Lead AI Architect — Feature engineering, XGBoost predictive modeling, and spatial graph propagation algorithm.", COLOR_CYAN),
        ("SYSTEMS & TELEMETRY", "[ Team Member 2 ]", "Systems Engineer — Real-time telemetry ingestion pipeline, OS hardware API bindings, and runtime engine.", COLOR_GREEN),
        ("PRODUCT & OPERATIONS", "[ Team Member 3 ]", "Product Lead — Mission Control UX, DCIM integration strategy, and Yuva Yodha Challenge submission.", COLOR_ORANGE)
    ]
    t_w = Inches(3.7)
    for idx, (role_tag, member_name, desc, color) in enumerate(team_members):
        left_x = Inches(0.8) + idx * Inches(4.0)
        add_card(s10, left_x, Inches(2.1), t_w, Inches(2.6), color)
        tb_t = s10.shapes.add_textbox(left_x + Inches(0.2), Inches(2.25), t_w - Inches(0.4), Inches(2.3))
        tf_t = tb_t.text_frame
        tf_t.word_wrap = True
        p = tf_t.paragraphs[0]
        p.text = role_tag
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        p2 = tf_t.add_paragraph()
        p2.text = member_name
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT_MAIN
        p2.space_before = Pt(4)
        p3 = tf_t.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = COLOR_TEXT_MUTED
        p3.space_before = Pt(8)

    # Closing Callout
    add_card(s10, Inches(0.8), Inches(4.9), Inches(11.7), Inches(1.7), COLOR_CYAN)
    tb10_call = s10.shapes.add_textbox(Inches(1.0), Inches(5.05), Inches(11.3), Inches(1.4))
    tf10_call = tb10_call.text_frame
    tf10_call.word_wrap = True
    p = tf10_call.paragraphs[0]
    p.text = "\"Every computation creates heat. We believe cooling should be intelligent enough to know where that heat is going next.\""
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.alignment = PP_ALIGN.CENTER

    p2 = tf10_call.add_paragraph()
    p2.text = "Thervo — \"Predict before you cool.\""
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_CYAN
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)

    add_footer(s10, 10)

    output_path = os.path.join(r"C:\Claude_projects\cooling project", "thervo_pitch_deck.pptx")
    prs.save(output_path)
    print(f"Presentation successfully created at: {output_path}")

if __name__ == "__main__":
    create_deck()
